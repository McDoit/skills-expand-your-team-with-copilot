"""
create_presentation.py
======================
Generates "agent-work-in-github-copilot.pptx" — a 4-slide, ~5-minute presentation
about the GitHub Copilot Coding Agent, based on the skills/expand-your-team-with-copilot exercise.

Usage:
    pip install python-pptx
    python create_presentation.py

The file "agent-work-in-github-copilot.pptx" will be saved in the current directory.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ---------------------------------------------------------------------------
# Theme colors
# ---------------------------------------------------------------------------
BG_COLOR = RGBColor(0x1B, 0x1F, 0x23)      # GitHub dark background
TITLE_COLOR = RGBColor(0x58, 0xA6, 0xFF)    # GitHub blue
ACCENT_COLOR = RGBColor(0x3F, 0xB9, 0x50)  # GitHub green
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREY = RGBColor(0x8B, 0x94, 0x9E)           # subtle text

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def set_bg(slide, color):
    """Fill slide background with a solid color."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, text, left, top, width, height,
                font_size=18, bold=False, color=WHITE,
                align=PP_ALIGN.LEFT, italic=False):
    """Add a simple textbox and return the text frame."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return tf


def add_bullet_paragraph(tf, text, level=0, font_size=16,
                          color=WHITE, bold=False, accent=False):
    """Append a bullet paragraph to an existing text frame."""
    p = tf.add_paragraph()
    p.level = level
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.name = "Calibri"
    run.font.color.rgb = ACCENT_COLOR if accent else color
    return p


def add_speaker_note(slide, note_text):
    """Add speaker notes to a slide."""
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = note_text


def add_slide_title(slide, title_text):
    """Add a large title at the top of the slide."""
    tf = add_textbox(
        slide, title_text,
        left=Inches(0.4), top=Inches(0.25),
        width=Inches(12.5), height=Inches(0.9),
        font_size=32, bold=True, color=TITLE_COLOR,
        align=PP_ALIGN.LEFT
    )
    return tf


def add_divider(slide, top):
    """Add a thin colored line as a visual divider."""
    from pptx.util import Emu
    line = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(0.4), top, Inches(12.5), Inches(0.04)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = TITLE_COLOR
    line.line.fill.background()


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------

def build_slide1(prs):
    """Slide 1 — What Is the Copilot Coding Agent?"""
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)
    set_bg(slide, BG_COLOR)

    # Title
    add_slide_title(slide, "🤖  What Is the Copilot Coding Agent?")
    add_divider(slide, Inches(1.2))

    # Subtitle
    add_textbox(
        slide,
        "From code editor to autopilot: Copilot now works directly on GitHub — no IDE required.",
        left=Inches(0.4), top=Inches(1.3),
        width=Inches(12.5), height=Inches(0.5),
        font_size=16, italic=True, color=GREY
    )

    # Comparison table header
    tf = add_textbox(
        slide, "Copilot in the Editor  vs.  Copilot Coding Agent",
        left=Inches(0.4), top=Inches(1.9),
        width=Inches(12.5), height=Inches(0.4),
        font_size=17, bold=True, color=ACCENT_COLOR
    )

    # Comparison rows as bullet points
    rows = [
        ("Interface",   "Code editor",              "Issues & Pull Requests"),
        ("Scope",       "Local files",              "Entire repository"),
        ("Activation",  "Chat / inline suggestions","Assign issue to @copilot"),
        ("MCP Support", "✅",                       "✅"),
    ]
    bullets_tf = add_textbox(
        slide, "",
        left=Inches(0.4), top=Inches(2.35),
        width=Inches(12.5), height=Inches(2.0),
        font_size=15
    )
    first = True
    for label, editor_val, agent_val in rows:
        p = bullets_tf.paragraphs[0] if first else bullets_tf.add_paragraph()
        first = False
        run = p.add_run()
        run.text = f"  {label:<16}  Editor: {editor_val:<36}  Agent: {agent_val}"
        run.font.size = Pt(15)
        run.font.name = "Calibri"
        run.font.color.rgb = WHITE

    # Key point
    kp_tf = add_textbox(
        slide, "Key Point",
        left=Inches(0.4), top=Inches(4.4),
        width=Inches(2.0), height=Inches(0.4),
        font_size=16, bold=True, color=ACCENT_COLOR
    )
    add_textbox(
        slide,
        "Anyone with write access can assign a GitHub Issue to Copilot. "
        "It creates a branch, opens a PR, implements code, and requests your review — "
        "like a human teammate.",
        left=Inches(0.4), top=Inches(4.8),
        width=Inches(12.5), height=Inches(0.9),
        font_size=16, color=WHITE
    )

    add_speaker_note(
        slide,
        "Think of it as adding a junior developer to your team who never sleeps. "
        "You write a clear issue, assign it, and Copilot handles the coding while "
        "you focus on other work."
    )


def build_slide2(prs):
    """Slide 2 — The Workflow: Assign → Review → Merge"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_bg(slide, BG_COLOR)

    add_slide_title(slide, "The Workflow: Assign → Review → Merge")
    add_divider(slide, Inches(1.2))

    steps = [
        "1.  📝  Write a clear issue — with goal & acceptance criteria",
        "2.  👤  Assign to Copilot — select Copilot in the Assignees dropdown",
        "3.  🤖  Copilot creates a branch & PR — works inside a GitHub Actions session",
        "4.  👀  Watch progress live — session logs stream in real-time",
        "5.  💬  Give feedback — use @copilot in PR review comments",
        "6.  🔄  Iterate — Copilot implements feedback, you re-review",
        "7.  ✅  Merge — when satisfied, merge just like any other PR",
    ]

    tf = add_textbox(
        slide, steps[0],
        left=Inches(0.5), top=Inches(1.35),
        width=Inches(8.0), height=Inches(3.8),
        font_size=15, color=WHITE
    )
    for step in steps[1:]:
        add_bullet_paragraph(tf, step, font_size=15)

    # Demo scenario box
    add_textbox(
        slide, "Demo scenario",
        left=Inches(8.8), top=Inches(1.35),
        width=Inches(4.1), height=Inches(0.35),
        font_size=14, bold=True, color=ACCENT_COLOR
    )
    add_textbox(
        slide,
        'A teacher asks to add a "Manga Maniacs" club to the school website. '
        "They file an issue, Copilot writes the code, the teacher reviews and asks "
        "for a more manga-inspired description — Copilot revises. All without opening VS Code.",
        left=Inches(8.8), top=Inches(1.75),
        width=Inches(4.1), height=Inches(2.0),
        font_size=13, color=GREY
    )

    add_speaker_note(
        slide,
        "Walk through the assign → PR → review → merge flow. "
        "Emphasize that @copilot in comments triggers new sessions, "
        "while regular comments (without @copilot) are ignored by the agent."
    )


def build_slide3(prs):
    """Slide 3 — Customizing & Scaling with the Agents Panel"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_bg(slide, BG_COLOR)

    add_slide_title(slide, "Customizing & Scaling with the Agents Panel")
    add_divider(slide, Inches(1.2))

    # Customization files
    add_textbox(
        slide, "Two Key Customization Files",
        left=Inches(0.4), top=Inches(1.35),
        width=Inches(6.0), height=Inches(0.4),
        font_size=17, bold=True, color=ACCENT_COLOR
    )

    files_tf = add_textbox(
        slide,
        ".github/copilot-instructions.md",
        left=Inches(0.4), top=Inches(1.8),
        width=Inches(6.0), height=Inches(1.8),
        font_size=15, bold=True, color=TITLE_COLOR
    )
    add_bullet_paragraph(
        files_tf,
        "  Repo-specific context: architecture, conventions, standards",
        font_size=14, color=WHITE
    )
    add_bullet_paragraph(
        files_tf,
        ".github/workflows/copilot-setup-steps.yml",
        font_size=15, bold=True, color=TITLE_COLOR
    )
    add_bullet_paragraph(
        files_tf,
        "  Pre-install tools & dependencies for the agent's environment",
        font_size=14, color=WHITE
    )

    # Agents panel
    add_textbox(
        slide, "Agents Panel  —  github.com/copilot/agents",
        left=Inches(6.8), top=Inches(1.35),
        width=Inches(6.1), height=Inches(0.4),
        font_size=17, bold=True, color=ACCENT_COLOR
    )

    panel_tf = add_textbox(
        slide,
        "🛠️  Assign tasks without leaving your current page",
        left=Inches(6.8), top=Inches(1.8),
        width=Inches(6.1), height=Inches(1.8),
        font_size=15, color=WHITE
    )
    add_bullet_paragraph(
        panel_tf,
        "👀  Monitor multiple parallel tasks with real-time status",
        font_size=15, color=WHITE
    )
    add_bullet_paragraph(
        panel_tf,
        "🔗  Jump directly to the PR when ready to review",
        font_size=15, color=WHITE
    )

    # Exercise note
    add_textbox(
        slide,
        "In the exercise, learners assign 3 tasks simultaneously "
        "(social sharing, difficulty tracks, dark mode) and track them all from one panel.",
        left=Inches(0.4), top=Inches(3.7),
        width=Inches(12.5), height=Inches(0.7),
        font_size=15, italic=True, color=GREY
    )

    add_speaker_note(
        slide,
        "Custom instructions = onboarding docs for your AI teammate. "
        "Setup workflow = pre-configured dev machine. "
        "Agents Panel = mission control."
    )


def build_slide4(prs):
    """Slide 4 — Safety & Takeaways"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_bg(slide, BG_COLOR)

    add_slide_title(slide, "Safety & Takeaways")
    add_divider(slide, Inches(1.2))

    # Guardrails
    add_textbox(
        slide, "Built-in Guardrails",
        left=Inches(0.4), top=Inches(1.35),
        width=Inches(6.2), height=Inches(0.4),
        font_size=17, bold=True, color=ACCENT_COLOR
    )

    guardrails = [
        "🔒  Agent only modifies its own copilot/* branch",
        "🔥  Configurable firewall restricts internet access",
        "🛡️  Branch protections & rulesets still enforced",
        "👁️  Only write-access users can assign; hidden issue content ignored",
        "👤  Commits co-authored by assigner (contribution graph safe)",
    ]
    gr_tf = add_textbox(
        slide, guardrails[0],
        left=Inches(0.4), top=Inches(1.8),
        width=Inches(6.2), height=Inches(2.5),
        font_size=14, color=WHITE
    )
    for g in guardrails[1:]:
        add_bullet_paragraph(gr_tf, g, font_size=14, color=WHITE)

    # Takeaways
    add_textbox(
        slide, "Key Takeaways",
        left=Inches(6.8), top=Inches(1.35),
        width=Inches(6.1), height=Inches(0.4),
        font_size=17, bold=True, color=ACCENT_COLOR
    )

    takeaways = [
        "1.  Copilot coding agent turns issues into PRs — no editor needed",
        "2.  Clear issues = better results (treat it like onboarding a teammate)",
        "3.  Custom instructions + setup steps improve quality and speed",
        "4.  Parallel task delegation via Agents Panel multiplies throughput",
    ]
    ta_tf = add_textbox(
        slide, takeaways[0],
        left=Inches(6.8), top=Inches(1.8),
        width=Inches(6.1), height=Inches(2.5),
        font_size=14, color=WHITE
    )
    for t in takeaways[1:]:
        add_bullet_paragraph(ta_tf, t, font_size=14, color=WHITE)

    # Try it link
    add_textbox(
        slide,
        "💡  Try it yourself:  github.com/skills/expand-your-team-with-copilot",
        left=Inches(0.4), top=Inches(4.5),
        width=Inches(12.5), height=Inches(0.45),
        font_size=16, bold=True, color=TITLE_COLOR
    )

    add_speaker_note(
        slide,
        "Close by reinforcing that this isn't uncontrolled AI — every change goes through "
        "standard code review. The exercise takes less than 1 hour and gives hands-on "
        "experience with the full agent loop."
    )


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    build_slide1(prs)
    build_slide2(prs)
    build_slide3(prs)
    build_slide4(prs)

    output_path = "agent-work-in-github-copilot.pptx"
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")


if __name__ == "__main__":
    main()
